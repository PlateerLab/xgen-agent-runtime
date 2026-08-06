"""Browser* (an-web) + Doc* (edit2docs) built-in tool families (2.43.0).

The engines are optional extras — every test that needs one skips
cleanly when it is not importable, so the suite stays green on minimal
installs. Deterministic paths run against the real engines with local
fixtures (no network, no LLM key).
"""

from __future__ import annotations

import json
import os
import sys
from pathlib import Path

import pytest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "..", "src"))

from xgen_agent_runtime.tools.base import ToolContext
from xgen_agent_runtime.tools.built_in import BUILT_IN_TOOL_CLASSES, BUILT_IN_TOOL_FEATURES
from xgen_agent_runtime.tools.built_in.browser_tools import (
    BROWSER_TOOL_CLASSES,
    BrowserActTool,
    BrowserCloseTool,
    BrowserNavigateTool,
    BrowserSnapshotTool,
    _parse_target,
    _runtime,
)
from xgen_agent_runtime.tools.built_in.doc_tools import (
    DOC_TOOL_CLASSES,
    DocAnalyzeTool,
    DocApplyEditsTool,
    DocBuildTool,
    DocEditTool,
    DocGenerateTool,
    DocGuideTool,
    DocRenderTool,
    DocXmlEditTool,
    DocXmlReadTool,
)

an_web = pytest.importorskip("an_web", reason="an-web extra not installed")
edit2docs = pytest.importorskip("edit2docs", reason="edit2docs extra not installed")


# ── Registration ───────────────────────────────────────────


class TestRegistration:
    def test_families_registered(self):
        for name in list(BROWSER_TOOL_CLASSES) + list(DOC_TOOL_CLASSES):
            assert name in BUILT_IN_TOOL_CLASSES

    def test_feature_groups(self):
        assert BUILT_IN_TOOL_FEATURES["browser"] == list(BROWSER_TOOL_CLASSES.keys())
        assert BUILT_IN_TOOL_FEATURES["documents"] == list(DOC_TOOL_CLASSES.keys())

    def test_schemas_are_valid_shapes(self):
        for name, cls in {**BROWSER_TOOL_CLASSES, **DOC_TOOL_CLASSES}.items():
            tool = cls()
            fmt = tool.to_api_format()
            assert fmt["name"] == name
            assert fmt["description"]
            assert fmt["input_schema"]["type"] == "object"


# ── Browser target parsing ─────────────────────────────────


class TestTargetParsing:
    def test_ref_handle(self):
        assert _parse_target("n42") == {"by": "node_id", "node_id": "n42"}

    def test_text_prefix(self):
        assert _parse_target("text=Sign in") == {"by": "text", "text": "Sign in"}

    def test_css_passthrough(self):
        assert _parse_target("#login .btn") == "#login .btn"

    def test_dict_passthrough(self):
        loc = {"by": "role", "role": "button", "text": "Go"}
        assert _parse_target(loc) is loc


# ── Browser tools against a local HTML page (no real network) ──


@pytest.fixture
def html_url(tmp_path, monkeypatch):
    """Serve a small page over HTTP from localhost (an-web fetches for
    real; a loopback server keeps the test hermetic)."""
    import http.server
    import threading

    # The 2.51.1 SSRF guard blocks loopback by default; this hermetic
    # fixture legitimately targets 127.0.0.1, so opt into the escape hatch.
    monkeypatch.setenv("GENY_ALLOW_PRIVATE_URLS", "1")

    page = (
        "<html><head><title>Fixture Page</title></head><body>"
        "<h1>Hello World</h1>"
        "<p id='intro'>Intro text with <a href='/next.html' id='go'>a link</a>.</p>"
        "<button id='btn'>Press me</button>"
        "</body></html>"
    )
    next_page = (
        "<html><head><title>Next Page</title></head><body>"
        "<h1>Second page</h1><p>You arrived.</p></body></html>"
    )
    (tmp_path / "index.html").write_text(page, encoding="utf-8")
    (tmp_path / "next.html").write_text(next_page, encoding="utf-8")

    handler = lambda *a, **kw: http.server.SimpleHTTPRequestHandler(  # noqa: E731
        *a, directory=str(tmp_path), **kw
    )
    server = http.server.ThreadingHTTPServer(("127.0.0.1", 0), handler)
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()
    try:
        yield f"http://127.0.0.1:{server.server_address[1]}/index.html"
    finally:
        server.shutdown()
        thread.join(timeout=5)


class TestBrowserFlow:
    @pytest.mark.asyncio
    async def test_navigate_snapshot_click_close(self, html_url):
        sid = "browser-flow-test"
        ctx = ToolContext(session_id=sid)
        try:
            nav = await BrowserNavigateTool().execute({"url": html_url}, ctx)
            assert not nav.is_error, nav.content
            assert "Fixture Page" in nav.content
            assert "Hello World" in nav.content
            assert "[ref=" in nav.content  # interactive elements got handles

            snap = await BrowserSnapshotTool().execute({}, ctx)
            assert not snap.is_error
            assert "Fixture Page" in snap.content

            # Click the link — navigates to next.html and inlines the new page.
            act = await BrowserActTool().execute(
                {"action": "click", "target": "#go"}, ctx
            )
            assert not act.is_error, act.content
            assert "Second page" in act.content or "Next Page" in act.content

            closed = await BrowserCloseTool().execute({}, ctx)
            assert not closed.is_error
            assert closed.metadata["closed"] is True
        finally:
            await _runtime.close_session(sid)

    @pytest.mark.asyncio
    async def test_act_without_page_errors(self):
        ctx = ToolContext(session_id="browser-no-page")
        result = await BrowserActTool().execute(
            {"action": "click", "target": "#x"}, ctx
        )
        assert result.is_error
        assert "BrowserNavigate" in result.content

    @pytest.mark.asyncio
    async def test_navigate_rejects_bad_scheme(self):
        ctx = ToolContext(session_id="browser-bad-scheme")
        result = await BrowserNavigateTool().execute({"url": "file:///etc/passwd"}, ctx)
        assert result.is_error

    @pytest.mark.asyncio
    async def test_sessions_are_isolated(self, html_url):
        ctx_a = ToolContext(session_id="browser-iso-a")
        ctx_b = ToolContext(session_id="browser-iso-b")
        try:
            await BrowserNavigateTool().execute({"url": html_url}, ctx_a)
            # Session B never navigated — it must not see A's tab.
            snap_b = await BrowserSnapshotTool().execute({}, ctx_b)
            assert snap_b.is_error
        finally:
            await _runtime.close_session("browser-iso-a")
            await _runtime.close_session("browser-iso-b")


# ── Doc tools against real generated fixtures (no LLM) ─────


@pytest.fixture
def docx_path(tmp_path):
    from edit2docs.documents.docx_engine import docx_from_markdown

    data = docx_from_markdown("# Title\n\nFirst paragraph.\n\nSecond paragraph.")
    p = tmp_path / "doc.docx"
    p.write_bytes(data)
    return p


@pytest.fixture
def xlsx_path(tmp_path):
    from edit2docs.documents.xlsx_engine import xlsx_from_spec

    data = xlsx_from_spec(
        {"sheets": [{"name": "Data", "headers": ["a", "b"], "rows": [[1, 2], [3, 4]]}]}
    )
    p = tmp_path / "book.xlsx"
    p.write_bytes(data)
    return p


@pytest.fixture
def chart_pptx_path(tmp_path):
    pptx = pytest.importorskip("pptx", reason="python-pptx not installed")
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    cd = CategoryChartData()
    cd.categories = ["A", "B", "C"]
    cd.add_series("S1", (1, 2, 3))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(6), Inches(4), cd
    )
    p = tmp_path / "deck.pptx"
    prs.save(str(p))
    return p


class TestDocTools:
    @pytest.mark.asyncio
    async def test_analyze_docx(self, docx_path, tmp_path):
        ctx = ToolContext(working_dir=str(tmp_path))
        result = await DocAnalyzeTool().execute({"path": "doc.docx"}, ctx)
        assert not result.is_error, result.content
        info = json.loads(result.content)
        assert info["format"] == "docx"
        assert any("para" in item for item in info["outline"])

    @pytest.mark.asyncio
    async def test_apply_edits_xlsx_roundtrip(self, xlsx_path, tmp_path):
        ctx = ToolContext(working_dir=str(tmp_path))
        result = await DocApplyEditsTool().execute(
            {
                "path": "book.xlsx",
                "edits": [{"action": "set_cell", "sheet": "Data", "cell": "B2", "value": 99}],
            },
            ctx,
        )
        assert not result.is_error, result.content
        summary = json.loads(result.content)
        assert summary["applied"] == 1
        assert summary["failed"] == 0

        # The edit is visible to a fresh analyze (in-place default output).
        check = await DocAnalyzeTool().execute({"path": summary["path"]}, ctx)
        assert "99" in check.content

    @pytest.mark.asyncio
    async def test_edit_chart_retitle_and_data(self, chart_pptx_path, tmp_path):
        """Chart edits ride DocApplyEdits — a `chart` key routes them."""
        ctx = ToolContext(working_dir=str(tmp_path))
        result = await DocApplyEditsTool().execute(
            {
                "path": "deck.pptx",
                "edits": [
                    {"chart": 0, "title": "Q3 Sales"},
                    {
                        "chart": 0,
                        "categories": ["Q1", "Q2", "Q3"],
                        "series": [{"name": "Rev", "values": [10, 20, 30]}],
                    },
                ],
            },
            ctx,
        )
        assert not result.is_error, result.content
        summary = json.loads(result.content)
        assert summary["applied"] == 2
        assert summary["failed"] == 0
        # The retitle + data change is visible to a fresh analyze.
        check = await DocAnalyzeTool().execute({"path": summary["path"]}, ctx)
        assert "Q3 Sales" in check.content
        assert "Rev" in check.content

    @pytest.mark.asyncio
    async def test_build_docx_from_markdown(self, tmp_path):
        ctx = ToolContext(working_dir=str(tmp_path))
        result = await DocBuildTool().execute(
            {"spec": "# Report\n\nBody **text**.\n\n- a\n- b", "output": "out.docx"},
            ctx,
        )
        assert not result.is_error, result.content
        summary = json.loads(result.content)
        assert (tmp_path / "out.docx").exists()
        # Round-trips through analyze.
        check = await DocAnalyzeTool().execute({"path": summary["path"]}, ctx)
        assert json.loads(check.content)["format"] == "docx"

    @pytest.mark.asyncio
    async def test_build_pptx_from_slide_spec(self, tmp_path):
        ctx = ToolContext(working_dir=str(tmp_path))
        result = await DocBuildTool().execute(
            {
                "spec": {"slides": [
                    {"layout": "title", "title": "Deck", "subtitle": "2026"},
                    {"layout": "content", "title": "Agenda", "bullets": ["A", "B"]},
                ]},
                "output": "deck.pptx",
            },
            ctx,
        )
        assert not result.is_error, result.content
        assert json.loads(result.content)["page_count"] == 2
        assert (tmp_path / "deck.pptx").exists()

    @pytest.mark.asyncio
    async def test_xml_read_lists_parts_and_reads_one(self, chart_pptx_path, tmp_path):
        ctx = ToolContext(working_dir=str(tmp_path))
        listing = await DocXmlReadTool().execute({"path": "deck.pptx"}, ctx)
        assert not listing.is_error, listing.content
        parts = json.loads(listing.content)["parts"]
        assert any("charts/chart1.xml" in p["part"] for p in parts)
        read = await DocXmlReadTool().execute(
            {"path": "deck.pptx", "part": "ppt/charts/chart1.xml"}, ctx
        )
        assert not read.is_error and "<c:ser>" in read.content

    @pytest.mark.asyncio
    async def test_xml_edit_recolors_chart_series(self, chart_pptx_path, tmp_path):
        """The real-world failure case: recolor bars — now a pure tool call."""
        pptx = pytest.importorskip("pptx")
        ctx = ToolContext(working_dir=str(tmp_path))
        result = await DocXmlEditTool().execute(
            {
                "path": "deck.pptx",
                "part": "ppt/charts/chart1.xml",
                "edits": [{
                    "find": "</c:tx>",
                    "replace": (
                        "</c:tx><c:spPr><a:solidFill>"
                        '<a:srgbClr val="FF0000"/>'
                        "</a:solidFill></c:spPr>"
                    ),
                }],
            },
            ctx,
        )
        assert not result.is_error, result.content
        assert json.loads(result.content)["applied"] == 1
        prs = pptx.Presentation(str(chart_pptx_path))
        chart = next(s for sl in prs.slides for s in sl.shapes if s.has_chart).chart
        assert str(chart.series[0].format.fill.fore_color.rgb) == "FF0000"

    @pytest.mark.asyncio
    async def test_xml_edit_rejects_malformed_result(self, chart_pptx_path, tmp_path):
        ctx = ToolContext(working_dir=str(tmp_path))
        result = await DocXmlEditTool().execute(
            {
                "path": "deck.pptx",
                "part": "ppt/charts/chart1.xml",
                "edits": [{"find": "</c:chartSpace>", "replace": "<broken"}],
            },
            ctx,
        )
        # Refused as engine feedback: nothing applied, doc still valid.
        assert not result.is_error
        summary = json.loads(result.content)
        assert summary["applied"] == 0 and summary["failed"] >= 1

    @pytest.mark.asyncio
    async def test_guide_root_map_with_executor_names(self, tmp_path):
        """DocGuide is the skill entry point: family map rendered with the
        EXECUTOR tool names (never the library verb names)."""
        ctx = ToolContext(working_dir=str(tmp_path))
        result = await DocGuideTool().execute({}, ctx)
        assert not result.is_error, result.content
        assert "GENERATE" in result.content and "EDIT" in result.content
        assert "DocApplyEdits" in result.content
        assert "DocXmlEdit" in result.content
        assert "set_doc_text" not in result.content
        assert result.metadata["topics"]

    @pytest.mark.asyncio
    async def test_guide_topic_and_prefix(self, tmp_path):
        ctx = ToolContext(working_dir=str(tmp_path))
        colors = await DocGuideTool().execute({"topic": "recipes.colors"}, ctx)
        assert "srgbClr" in colors.content and "DocXmlEdit" in colors.content
        recipes = await DocGuideTool().execute({"topic": "recipes"}, ctx)
        assert "ADD A SLIDE" in recipes.content and "RECOLOR" in recipes.content
        unknown = await DocGuideTool().execute({"topic": "zzz"}, ctx)
        assert not unknown.is_error and "GENERATE" in unknown.content

    def test_descriptions_stay_compact(self):
        """Progressive-disclosure contract: frontmatter tier stays small;
        the fat how-to lives behind DocGuide(topic)."""
        for name, cls in DOC_TOOL_CLASSES.items():
            desc = cls().description
            assert len(desc) <= 320, (
                f"{name} description grew to {len(desc)} chars — move "
                "detail into edit2docs agent_guide GUIDES instead"
            )

    def test_guide_registered_first(self):
        assert list(DOC_TOOL_CLASSES)[0] == "DocGuide"

    def test_llm_verbs_are_feature_gated(self):
        """Keyless hosts must never see DocGenerate/DocEdit — they advertise
        feature:docs_llm so progressive disclosure drops them."""
        assert DocGenerateTool().required_config_keys() == ["feature:docs_llm"]
        assert DocEditTool().required_config_keys() == ["feature:docs_llm"]

    @pytest.mark.asyncio
    async def test_xml_edit_creates_and_deletes_parts(self, chart_pptx_path, tmp_path):
        ctx = ToolContext(working_dir=str(tmp_path))
        # create a brand-new XML part (content type registered)
        created = await DocXmlEditTool().execute(
            {
                "path": "deck.pptx",
                "part": "ppt/slides/slide2.xml",
                "xml": (
                    await DocXmlReadTool().execute(
                        {"path": "deck.pptx", "part": "ppt/slides/slide1.xml"}, ctx
                    )
                ).content,
                "content_type": (
                    "application/vnd.openxmlformats-officedocument"
                    ".presentationml.slide+xml"
                ),
            },
            ctx,
        )
        assert not created.is_error, created.content
        listing = await DocXmlReadTool().execute({"path": "deck.pptx"}, ctx)
        names = [q["part"] for q in json.loads(listing.content)["parts"]]
        assert "ppt/slides/slide2.xml" in names
        # delete it again
        deleted = await DocXmlEditTool().execute(
            {"path": "deck.pptx", "part": "ppt/slides/slide2.xml", "delete": True},
            ctx,
        )
        assert not deleted.is_error, deleted.content
        listing = await DocXmlReadTool().execute({"path": "deck.pptx"}, ctx)
        names = [q["part"] for q in json.loads(listing.content)["parts"]]
        assert "ppt/slides/slide2.xml" not in names

    @pytest.mark.asyncio
    async def test_xml_edit_requires_exactly_one_mode(self, chart_pptx_path, tmp_path):
        ctx = ToolContext(working_dir=str(tmp_path))
        result = await DocXmlEditTool().execute(
            {"path": "deck.pptx", "part": "ppt/charts/chart1.xml"}, ctx
        )
        assert result.is_error

    @pytest.mark.asyncio
    async def test_build_rejects_wrong_spec_type(self, tmp_path):
        ctx = ToolContext(working_dir=str(tmp_path))
        result = await DocBuildTool().execute(
            {"spec": "markdown-not-a-dict", "output": "x.pptx"}, ctx
        )
        assert result.is_error

    @pytest.mark.asyncio
    async def test_edit_chart_out_of_range_soft_fails(self, chart_pptx_path, tmp_path):
        ctx = ToolContext(working_dir=str(tmp_path))
        result = await DocApplyEditsTool().execute(
            {"path": "deck.pptx", "edits": [{"chart": 9, "title": "nope"}]},
            ctx,
        )
        # Out-of-range is soft engine feedback, not a hard tool error.
        assert not result.is_error, result.content
        summary = json.loads(result.content)
        assert summary["applied"] == 0
        assert summary["failed"] == 1
        assert summary["results"][0]["status"] == "not_found"

    @pytest.mark.asyncio
    async def test_apply_edits_soft_fail_reports_status(self, xlsx_path, tmp_path):
        ctx = ToolContext(working_dir=str(tmp_path))
        result = await DocApplyEditsTool().execute(
            {
                "path": "book.xlsx",
                "edits": [{"action": "set_cell", "sheet": "Nope", "cell": "A1", "value": 1}],
            },
            ctx,
        )
        assert not result.is_error  # soft-fail: statuses, not exceptions
        summary = json.loads(result.content)
        assert summary["applied"] == 0
        assert summary["failed"] == 1

    @pytest.mark.asyncio
    async def test_render_md_readable_content(self, docx_path, tmp_path):
        """DocPreview folded into DocRender: to='md' returns readable content."""
        ctx = ToolContext(working_dir=str(tmp_path))
        result = await DocRenderTool().execute({"path": "doc.docx", "to": "md"}, ctx)
        assert not result.is_error, result.content
        payload = json.loads(result.content)
        assert payload["to"] == "md"
        md = Path(payload["paths"][0]).read_text(encoding="utf-8")
        assert "First paragraph" in md

    @pytest.mark.asyncio
    async def test_path_guard_blocks_escape(self, tmp_path):
        inner = tmp_path / "inner"
        inner.mkdir()
        ctx = ToolContext(working_dir=str(inner), allowed_paths=[str(inner)])
        result = await DocAnalyzeTool().execute({"path": "../../etc/passwd"}, ctx)
        assert result.is_error
        assert "Access denied" in result.content or "No such file" in result.content

    @pytest.mark.asyncio
    async def test_llm_verbs_require_key(self, docx_path, tmp_path, monkeypatch):
        monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
        ctx = ToolContext(working_dir=str(tmp_path))
        gen = await DocGenerateTool().execute(
            {"intent": "x", "output": "new.docx"}, ctx
        )
        assert gen.is_error
        assert "ANTHROPIC_API_KEY" in gen.content
        edit = await DocEditTool().execute(
            {"path": "doc.docx", "instruction": "x"}, ctx
        )
        assert edit.is_error
        assert "ANTHROPIC_API_KEY" in edit.content

    @pytest.mark.asyncio
    async def test_render_png_pages(self, docx_path, tmp_path):
        from xgen_agent_runtime.tools.built_in.doc_tools import DocRenderTool

        ctx = ToolContext(working_dir=str(tmp_path))
        result = await DocRenderTool().execute(
            {"path": "doc.docx", "to": "png", "out_dir": "prev"}, ctx
        )
        assert not result.is_error, result.content
        payload = json.loads(result.content)
        assert payload["page_count"] >= 1
        assert payload["paths"][0].endswith("page-1.png")
        assert (tmp_path / "prev" / "page-1.png").exists()

    @pytest.mark.asyncio
    async def test_render_pdf(self, docx_path, tmp_path):
        from xgen_agent_runtime.tools.built_in.doc_tools import DocRenderTool

        ctx = ToolContext(working_dir=str(tmp_path))
        result = await DocRenderTool().execute({"path": "doc.docx", "to": "pdf"}, ctx)
        assert not result.is_error, result.content
        payload = json.loads(result.content)
        assert payload["to"] == "pdf" and payload["paths"][0].endswith(".pdf")

    @pytest.mark.asyncio
    async def test_unsupported_extension_rejected(self, tmp_path):
        (tmp_path / "notes.txt").write_text("hi", encoding="utf-8")
        ctx = ToolContext(working_dir=str(tmp_path))
        result = await DocAnalyzeTool().execute({"path": "notes.txt"}, ctx)
        assert result.is_error
        assert "Unsupported document format" in result.content
